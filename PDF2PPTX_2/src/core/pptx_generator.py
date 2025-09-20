"""
PowerPoint生成エンジン
python-pptxを使用したPowerPoint(.pptx)ファイル生成機能
"""

from pptx import Presentation
from pptx.util import Inches, Mm, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

import os
import io
from pathlib import Path
from typing import List, Tuple, Optional, Dict, Any, Union
from dataclasses import dataclass
from enum import Enum
from PIL import Image
from loguru import logger

from utils.error_handler import ErrorHandler, ErrorType
from utils.config import ConfigManager


class SlideSize(Enum):
    """スライドサイズ定義"""
    A3_LANDSCAPE = "a3_landscape"  # A3横（420×297mm）
    A4_LANDSCAPE = "a4_landscape"  # A4横（297×210mm）
    A4_PORTRAIT = "a4_portrait"    # A4縦（210×297mm）
    WIDESCREEN = "widescreen"      # ワイドスクリーン（16:9）
    STANDARD = "standard"          # 標準（4:3）


class LabelPosition(Enum):
    """ラベル配置位置"""
    TOP_LEFT = "top_left"
    TOP_CENTER = "top_center"
    TOP_RIGHT = "top_right"
    BOTTOM_LEFT = "bottom_left"
    BOTTOM_CENTER = "bottom_center"
    BOTTOM_RIGHT = "bottom_right"
    NONE = "none"


@dataclass
class SlideSettings:
    """スライド設定"""
    size: SlideSize = SlideSize.A3_LANDSCAPE
    width_mm: float = 420.0
    height_mm: float = 297.0


@dataclass
class LabelStyle:
    """ラベルスタイル設定"""
    font_name: str = "メイリオ"
    font_size: int = 18
    font_color: str = "#FFFFFF"
    background_color: str = "#1976D2"
    position: LabelPosition = LabelPosition.BOTTOM_CENTER
    margin_mm: float = 5.0
    padding_mm: float = 3.0
    add_page_numbers: bool = True
    add_filename: bool = True
    custom_text: str = ""


@dataclass
class ImageSettings:
    """画像配置設定"""
    maintain_aspect_ratio: bool = True
    center_horizontally: bool = True
    center_vertically: bool = True
    max_width_percent: float = 90.0
    max_height_percent: float = 80.0
    margin_mm: float = 10.0


class PPTXGenerator:
    """PowerPoint生成エンジンクラス"""

    # スライドサイズのプリセット（mm単位）
    SLIDE_SIZE_PRESETS = {
        SlideSize.A3_LANDSCAPE: (420.0, 297.0),
        SlideSize.A4_LANDSCAPE: (297.0, 210.0),
        SlideSize.A4_PORTRAIT: (210.0, 297.0),
        SlideSize.WIDESCREEN: (355.6, 200.0),  # 16:9
        SlideSize.STANDARD: (254.0, 190.5)     # 4:3
    }

    def __init__(self,
                 config_manager: Optional[ConfigManager] = None,
                 error_handler: Optional[ErrorHandler] = None):
        """
        初期化

        Args:
            config_manager: 設定管理インスタンス
            error_handler: エラーハンドラーインスタンス
        """
        self.config_manager = config_manager or ConfigManager()
        self.error_handler = error_handler or ErrorHandler()

        # 設定を読み込み
        self._load_pptx_settings()

        logger.info("PowerPoint生成エンジンを初期化しました")

    def _load_pptx_settings(self) -> None:
        """設定ファイルからPowerPoint生成設定を読み込み"""
        try:
            config = self.config_manager.get_config("conversion.pdf_to_pptx")

            # スライド設定
            slide_size_config = config.get("slide_size", [420, 297])
            self.default_slide_settings = SlideSettings(
                size=SlideSize.A3_LANDSCAPE,
                width_mm=slide_size_config[0],
                height_mm=slide_size_config[1]
            )

            # ラベルスタイル設定
            label_position_str = config.get("label_position", "top_left")
            logger.info(f"設定から読み込んだlabel_position: {label_position_str}")

            self.default_label_style = LabelStyle(
                font_name=config.get("font_name", "メイリオ"),
                font_size=config.get("font_size", 18),
                font_color=config.get("text_color", "#FFFFFF"),
                background_color=config.get("background_color", "#1976D2"),
                position=LabelPosition(label_position_str),
                add_page_numbers=config.get("add_page_numbers", True),
                add_filename=config.get("add_filename", True),
                margin_mm=0.0  # 座標0,0配置のためマージンを0に設定
            )

            logger.info(f"作成されたLabelStyle.position: {self.default_label_style.position}")

            # 画像配置設定
            self.default_image_settings = ImageSettings(
                maintain_aspect_ratio=config.get("maintain_aspect_ratio", True),
                center_horizontally=config.get("center_horizontally", True),
                center_vertically=config.get("center_vertically", True),
                max_width_percent=config.get("max_width_percent", 90.0),
                max_height_percent=config.get("max_height_percent", 80.0)
            )

            logger.debug("PowerPoint生成設定を読み込みました")

        except Exception as e:
            logger.warning(f"設定読み込みエラー、デフォルト設定を使用: {e}")
            self.default_slide_settings = SlideSettings()
            self.default_label_style = LabelStyle()
            self.default_image_settings = ImageSettings()

    def create_presentation(self,
                          slide_settings: Optional[SlideSettings] = None) -> Presentation:
        """
        新しいプレゼンテーションを作成

        Args:
            slide_settings: スライド設定

        Returns:
            Presentationオブジェクト

        Raises:
            RuntimeError: プレゼンテーション作成エラー
        """
        try:
            settings = slide_settings or self.default_slide_settings

            # 新しいプレゼンテーションを作成
            prs = Presentation()

            # スライドサイズを設定
            self._set_slide_size(prs, settings)

            logger.info(f"プレゼンテーションを作成しました: {settings.width_mm}×{settings.height_mm}mm")
            return prs

        except Exception as e:
            self.error_handler.handle_error(
                error_type=ErrorType.PPTX_GENERATION,
                error=e,
                details={"slide_settings": str(settings)}
            )
            raise

    def _set_slide_size(self, prs: Presentation, settings: SlideSettings) -> None:
        """
        スライドサイズを設定

        Args:
            prs: Presentationオブジェクト
            settings: スライド設定
        """
        try:
            # mmをEMU（English Metric Units）に変換
            # 1mm = 36000 EMU
            width_emu = int(settings.width_mm * 36000)
            height_emu = int(settings.height_mm * 36000)

            # スライドサイズを設定
            prs.slide_width = width_emu
            prs.slide_height = height_emu

            logger.debug(f"スライドサイズを設定しました: {settings.width_mm}×{settings.height_mm}mm")

        except Exception as e:
            logger.error(f"スライドサイズ設定エラー: {e}")
            raise

    def add_image_slide(self,
                       prs: Presentation,
                       image_path: Path,
                       page_number: Optional[int] = None,
                       custom_title: Optional[str] = None,
                       image_settings: Optional[ImageSettings] = None,
                       label_style: Optional[LabelStyle] = None) -> None:
        """
        画像スライドを追加

        Args:
            prs: Presentationオブジェクト
            image_path: 画像ファイルパス
            page_number: ページ番号
            custom_title: カスタムタイトル
            image_settings: 画像配置設定
            label_style: ラベルスタイル設定

        Raises:
            FileNotFoundError: 画像ファイルが存在しない
            RuntimeError: スライド追加エラー
        """
        try:
            if not image_path.exists():
                raise FileNotFoundError(f"画像ファイルが見つかりません: {image_path}")

            img_settings = image_settings or self.default_image_settings
            lbl_style = label_style or self.default_label_style

            # 空白スライドを追加
            slide_layout = prs.slide_layouts[6]  # 空白レイアウト
            slide = prs.slides.add_slide(slide_layout)

            # 画像を追加
            self._add_image_to_slide(slide, image_path, img_settings, prs)

            # ラベルを追加
            if lbl_style.position != LabelPosition.NONE:
                self._add_label_to_slide(
                    slide, image_path, page_number, custom_title, lbl_style, prs
                )

            logger.debug(f"画像スライドを追加しました: {image_path}")

        except Exception as e:
            self.error_handler.handle_error(
                error_type=ErrorType.PPTX_GENERATION,
                error=e,
                details={"image_path": str(image_path)}
            )
            raise

    def _add_image_to_slide(self,
                           slide,
                           image_path: Path,
                           settings: ImageSettings,
                           presentation) -> None:
        """
        スライドに画像を追加

        Args:
            slide: スライドオブジェクト
            image_path: 画像ファイルパス
            settings: 画像配置設定
        """
        try:
            # 画像の元サイズを取得
            with Image.open(image_path) as img:
                img_width, img_height = img.size

            # スライドサイズを取得（EMU単位）
            # プレゼンテーションオブジェクトから直接取得
            slide_width_emu = presentation.slide_width
            slide_height_emu = presentation.slide_height

            # 利用可能なサイズを計算（マージンを考慮）
            margin_emu = Mm(settings.margin_mm)
            available_width = slide_width_emu - (margin_emu * 2)
            available_height = slide_height_emu - (margin_emu * 2)

            # 最大サイズを制限
            max_width = available_width * (settings.max_width_percent / 100.0)
            max_height = available_height * (settings.max_height_percent / 100.0)

            # アスペクト比を維持しながら画像サイズを計算
            if settings.maintain_aspect_ratio:
                # 画像のアスペクト比
                aspect_ratio = img_width / img_height

                # 幅に合わせる場合と高さに合わせる場合を比較
                width_based_height = max_width / aspect_ratio
                height_based_width = max_height * aspect_ratio

                if width_based_height <= max_height:
                    # 幅基準でサイズ決定
                    final_width = max_width
                    final_height = width_based_height
                else:
                    # 高さ基準でサイズ決定
                    final_width = height_based_width
                    final_height = max_height
            else:
                final_width = max_width
                final_height = max_height

            # 画像の配置位置を計算
            if settings.center_horizontally:
                left = (slide_width_emu - final_width) / 2
            else:
                left = margin_emu

            if settings.center_vertically:
                top = (slide_height_emu - final_height) / 2
            else:
                top = margin_emu

            # 画像を挿入
            slide.shapes.add_picture(
                str(image_path),
                left=int(left),
                top=int(top),
                width=int(final_width),
                height=int(final_height)
            )

            logger.debug(f"画像を配置しました: サイズ({int(final_width/36000)}×{int(final_height/36000)}mm)")

        except Exception as e:
            logger.error(f"画像配置エラー: {e}")
            raise

    def _add_label_to_slide(self,
                           slide,
                           image_path: Path,
                           page_number: Optional[int],
                           custom_title: Optional[str],
                           style: LabelStyle,
                           presentation) -> None:
        """
        スライドにラベルを追加

        Args:
            slide: スライドオブジェクト
            image_path: 画像ファイルパス
            page_number: ページ番号
            custom_title: カスタムタイトル
            style: ラベルスタイル
        """
        try:
            # ラベルテキストを構築
            label_text = self._build_label_text(image_path, page_number, custom_title, style)

            if not label_text.strip():
                return

            # ラベルの配置位置とサイズを計算
            label_box = self._calculate_label_box(slide, style, presentation)

            # 1. 四角形オブジェクトを追加（背景）
            from pptx.enum.shapes import MSO_SHAPE

            # デバッグ出力
            logger.info(f"ラベル配置座標: left={label_box['left']}, top={label_box['top']}, width={label_box['width']}, height={label_box['height']}")

            rectangle = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left=label_box["left"],
                top=label_box["top"],
                width=label_box["width"],
                height=label_box["height"]
            )

            # 実際に配置された座標を確認
            logger.info(f"実際の四角形座標: left={rectangle.left}, top={rectangle.top}")

            # 四角形の背景色を設定
            rect_fill = rectangle.fill
            rect_fill.solid()
            rect_fill.fore_color.rgb = RGBColor.from_string(style.background_color.replace("#", ""))

            # 四角形の枠線を設定
            rect_line = rectangle.line
            rect_line.color.rgb = RGBColor.from_string(style.background_color.replace("#", ""))
            rect_line.width = Pt(1)

            # 2. テキストオブジェクトを四角形の上に追加
            textbox = slide.shapes.add_textbox(
                left=label_box["left"],
                top=label_box["top"],
                width=label_box["width"],
                height=label_box["height"]
            )

            # テキストフレームを設定
            text_frame = textbox.text_frame
            text_frame.clear()
            # 上下中央配置のためのマージン調整
            text_frame.margin_left = Mm(2)  # 左マージン
            text_frame.margin_right = Mm(2)  # 右マージン
            text_frame.margin_top = 0       # 上マージンを0に
            text_frame.margin_bottom = 0    # 下マージンを0に
            text_frame.word_wrap = True

            # 垂直方向中央配置
            from pptx.enum.text import MSO_ANCHOR
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # テキストを設定
            p = text_frame.add_paragraph()
            p.text = label_text
            p.alignment = PP_ALIGN.LEFT  # 水平方向は左寄せ

            # フォントスタイルを設定
            font = p.font
            font.name = style.font_name
            font.size = Pt(style.font_size)
            font.color.rgb = RGBColor.from_string(style.font_color.replace("#", ""))

            # テキストボックスの背景を透明にする
            textbox_fill = textbox.fill
            textbox_fill.background()

            # テキストボックスの枠線を無効化
            textbox_line = textbox.line
            textbox_line.fill.background()

            logger.debug(f"ラベルを追加しました: {label_text[:50]}...")

        except Exception as e:
            logger.error(f"ラベル追加エラー: {e}")
            # ラベル追加失敗は致命的ではないため、ログのみ

    def _build_label_text(self,
                         image_path: Path,
                         page_number: Optional[int],
                         custom_title: Optional[str],
                         style: LabelStyle) -> str:
        """
        ラベルテキストを構築

        Args:
            image_path: 画像ファイルパス
            page_number: ページ番号
            custom_title: カスタムタイトル
            style: ラベルスタイル

        Returns:
            ラベルテキスト
        """
        text_parts = []

        # ファイル名と数値（ページ番号）を「ファイル名_数値」形式で結合
        if style.add_filename and style.add_page_numbers and page_number:
            filename = self._clean_filename(image_path.stem)
            text_parts.append(f"{filename}_{page_number}")
        elif style.add_filename:
            # ページ番号がない場合はファイル名のみ
            filename = self._clean_filename(image_path.stem)
            text_parts.append(filename)
        elif style.add_page_numbers and page_number:
            # ファイル名がない場合はページ番号のみ
            text_parts.append(str(page_number))

        # カスタムテキスト
        if style.custom_text:
            text_parts.append(style.custom_text)

        # カスタムタイトル
        if custom_title:
            text_parts.append(custom_title)

        return " | ".join(text_parts)

    def _clean_filename(self, filename: str) -> str:
        """
        ファイル名をクリーンアップ

        例:
        - 大漢和辞典序文_page_013_13 → 大漢和辞典序文_13
        - filename_page_001_1 → filename_1
        - simple_filename → simple_filename

        Args:
            filename: 元のファイル名

        Returns:
            クリーンアップされたファイル名
        """
        import re

        # .pdfを除去
        if filename.lower().endswith('.pdf'):
            filename = filename[:-4]

        # 様々なページ番号パターンを処理

        # パターン1: _page_XXX_YY → _YY
        # 例: 大漢和辞典序文_page_013_13 → 大漢和辞典序文_13
        pattern1 = r'_page_\d+_(\d+)$'
        match = re.search(pattern1, filename)
        if match:
            base_name = filename[:match.start()]
            page_num = match.group(1)
            return f"{base_name}_{page_num}"

        # パターン2: _page_XXX → ファイル名のみ（ページ番号削除）
        # 例: filename_page_013 → filename
        pattern2 = r'_page_\d+$'
        if re.search(pattern2, filename):
            filename = re.sub(pattern2, '', filename)

        # パターン3: 末尾の_数字_数字 → _最後の数字
        # 例: filename_001_13 → filename_13
        pattern3 = r'_\d+_(\d+)$'
        match = re.search(pattern3, filename)
        if match:
            base_name = filename[:match.start()]
            page_num = match.group(1)
            return f"{base_name}_{page_num}"

        return filename

    def _calculate_label_box(self, slide, style: LabelStyle, presentation) -> Dict[str, int]:
        """
        ラベルボックスの位置とサイズを計算

        Args:
            slide: スライドオブジェクト
            style: ラベルスタイル

        Returns:
            ラベルボックス情報（left, top, width, height）
        """
        # プレゼンテーションからサイズを取得
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        margin = Mm(style.margin_mm)

        # デバッグ: 現在のstyle.positionを確認
        logger.info(f"_calculate_label_box呼び出し - style.position: {style.position} (型: {type(style.position)})")
        logger.info(f"LabelPosition.TOP_LEFT: {LabelPosition.TOP_LEFT} (型: {type(LabelPosition.TOP_LEFT)})")
        logger.info(f"比較結果: {style.position == LabelPosition.TOP_LEFT}")
        logger.info(f"値比較: {style.position.value if hasattr(style.position, 'value') else 'N/A'} == {LabelPosition.TOP_LEFT.value}")

        # ラベルサイズ（左上配置用に最適化）
        # より確実な条件判定：enum比較と文字列値比較の両方
        is_top_left = (style.position == LabelPosition.TOP_LEFT or
                      (hasattr(style.position, 'value') and style.position.value == "top_left"))

        logger.info(f"is_top_left判定結果: {is_top_left}")

        if is_top_left:
            logger.info("TOP_LEFT条件に一致しました！")
            # 左上の場合はコンパクトなサイズに調整、座標0,0に配置
            label_width = Mm(80)  # 約8cm幅
            label_height = Mm(12)  # 約1.2cm高さ
            left = Mm(0)  # 座標0の地点（最左端）- EMU単位に変換
            top = Mm(0)   # 座標0の地点（最上端）- EMU単位に変換

            # デバッグ用情報出力
            logger.info(f"TOP_LEFT配置: left={left} ({left} EMU), top={top} ({top} EMU)")
            logger.info(f"label_width={label_width} ({label_width} EMU), label_height={label_height} ({label_height} EMU)")
        else:
            logger.info(f"TOP_LEFT以外の配置: {style.position}")
            # その他の位置の場合は従来通り
            label_width = slide_width - (margin * 2)
            label_height = Mm(15)

            if style.position in [LabelPosition.TOP_CENTER, LabelPosition.TOP_RIGHT]:
                top = margin
            else:  # BOTTOM_*
                top = slide_height - margin - label_height

            if style.position in [LabelPosition.TOP_RIGHT, LabelPosition.BOTTOM_RIGHT]:
                left = slide_width - margin - label_width
            else:  # CENTER
                left = margin

        return {
            "left": int(left),
            "top": int(top),
            "width": int(label_width),
            "height": int(label_height)
        }

    def add_multiple_images(self,
                          prs: Presentation,
                          image_paths: List[Path],
                          start_page_number: int = 1,
                          image_settings: Optional[ImageSettings] = None,
                          label_style: Optional[LabelStyle] = None,
                          progress_callback: Optional[callable] = None) -> None:
        """
        複数の画像をスライドに追加

        Args:
            prs: Presentationオブジェクト
            image_paths: 画像ファイルパスのリスト
            start_page_number: 開始ページ番号
            image_settings: 画像配置設定
            label_style: ラベルスタイル設定
            progress_callback: 進捗コールバック関数 callback(current, total)

        Raises:
            RuntimeError: スライド追加エラー
        """
        try:
            total_images = len(image_paths)
            completed = 0

            for i, image_path in enumerate(image_paths):
                try:
                    page_number = start_page_number + i
                    self.add_image_slide(
                        prs, image_path, page_number,
                        image_settings=image_settings,
                        label_style=label_style
                    )
                    completed += 1

                    if progress_callback:
                        progress_callback(completed, total_images)

                except Exception as e:
                    logger.error(f"画像 {image_path} のスライド追加に失敗: {e}")
                    continue

            logger.info(f"複数画像スライド追加完了: {completed}/{total_images} 枚")

        except Exception as e:
            self.error_handler.handle_error(
                error_type=ErrorType.PPTX_GENERATION,
                error=e,
                details={"total_images": total_images}
            )
            raise

    def save_presentation(self,
                         prs: Presentation,
                         output_path: Path,
                         overwrite: bool = False) -> bool:
        """
        プレゼンテーションをファイルに保存

        Args:
            prs: Presentationオブジェクト
            output_path: 出力ファイルパス
            overwrite: 既存ファイルを上書きするか

        Returns:
            保存成功時True

        Raises:
            FileExistsError: ファイルが既に存在する（overwrite=False時）
            RuntimeError: 保存エラー
        """
        try:
            # 出力ディレクトリを作成
            output_path.parent.mkdir(parents=True, exist_ok=True)

            # 既存ファイルチェック
            if output_path.exists() and not overwrite:
                raise FileExistsError(f"ファイルが既に存在します: {output_path}")

            # 保存実行
            prs.save(str(output_path))

            # ファイルサイズを確認
            file_size = output_path.stat().st_size
            file_size_mb = file_size / (1024 * 1024)

            logger.info(f"PowerPointファイルを保存しました: {output_path} ({file_size_mb:.1f}MB)")
            return True

        except Exception as e:
            self.error_handler.handle_error(
                error_type=ErrorType.PPTX_GENERATION,
                error=e,
                details={"output_path": str(output_path)}
            )
            raise

    def convert_images_to_pptx(self,
                              image_paths: List[Path],
                              output_path: Path,
                              title: Optional[str] = None,
                              slide_settings: Optional[SlideSettings] = None,
                              image_settings: Optional[ImageSettings] = None,
                              label_style: Optional[LabelStyle] = None,
                              progress_callback: Optional[callable] = None) -> bool:
        """
        画像ファイル群からPowerPointファイルを作成

        Args:
            image_paths: 画像ファイルパスのリスト
            output_path: 出力PowerPointファイルパス
            title: プレゼンテーションタイトル
            slide_settings: スライド設定
            image_settings: 画像配置設定
            label_style: ラベルスタイル設定
            progress_callback: 進捗コールバック関数 callback(current, total, status)

        Returns:
            変換成功時True

        Raises:
            ValueError: 入力ファイルが空
            RuntimeError: 変換エラー
        """
        try:
            if not image_paths:
                raise ValueError("画像ファイルが指定されていません")

            # 存在する画像ファイルのみフィルタリング
            valid_images = [path for path in image_paths if path.exists()]

            if not valid_images:
                raise ValueError("有効な画像ファイルが見つかりません")

            if progress_callback:
                progress_callback(0, len(valid_images), "プレゼンテーションを作成中...")

            # プレゼンテーション作成
            prs = self.create_presentation(slide_settings)

            # タイトルスライドを追加（オプション）- 表紙は不要なためコメントアウト
            # if title:
            #     self._add_title_slide(prs, title)

            if progress_callback:
                progress_callback(0, len(valid_images), "画像スライドを追加中...")

            # 画像スライドを追加
            self.add_multiple_images(
                prs, valid_images, 1, image_settings, label_style,
                progress_callback=lambda c, t: progress_callback(c, t, f"スライド追加中... ({c}/{t})")
            )

            if progress_callback:
                progress_callback(len(valid_images), len(valid_images), "ファイル保存中...")

            # ファイル保存
            self.save_presentation(prs, output_path, overwrite=True)

            if progress_callback:
                progress_callback(len(valid_images), len(valid_images), "変換完了")

            logger.info(f"PowerPoint変換完了: {len(valid_images)} 枚の画像を変換しました")
            return True

        except Exception as e:
            self.error_handler.handle_error(
                error_type=ErrorType.PPTX_GENERATION,
                error=e,
                details={"image_count": len(image_paths), "output_path": str(output_path)}
            )
            raise

    def _add_title_slide(self, prs: Presentation, title: str) -> None:
        """
        タイトルスライドを追加

        Args:
            prs: Presentationオブジェクト
            title: タイトル文字列
        """
        try:
            # タイトルスライドレイアウトを使用
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)

            # タイトルと副題を設定
            slide.shapes.title.text = title

            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = f"PDF2PPTX v3.0で生成"

            logger.debug(f"タイトルスライドを追加しました: {title}")

        except Exception as e:
            logger.warning(f"タイトルスライド追加エラー: {e}")

    def get_presentation_info(self, pptx_path: Path) -> Dict[str, Any]:
        """
        PowerPointファイルの情報を取得

        Args:
            pptx_path: PowerPointファイルパス

        Returns:
            プレゼンテーション情報辞書

        Raises:
            FileNotFoundError: ファイルが存在しない
            ValueError: PowerPointファイルでない
        """
        try:
            if not pptx_path.exists():
                raise FileNotFoundError(f"ファイルが見つかりません: {pptx_path}")

            prs = Presentation(str(pptx_path))

            file_stats = pptx_path.stat()

            info = {
                "file_path": str(pptx_path),
                "file_size": file_stats.st_size,
                "file_size_mb": file_stats.st_size / (1024 * 1024),
                "slide_count": len(prs.slides),
                "slide_width_mm": prs.slide_width / 36000,
                "slide_height_mm": prs.slide_height / 36000,
                "creation_time": file_stats.st_ctime,
                "modification_time": file_stats.st_mtime,
                "slides": []
            }

            # 各スライドの情報
            for i, slide in enumerate(prs.slides):
                slide_info = {
                    "slide_number": i + 1,
                    "layout_name": slide.slide_layout.name,
                    "shape_count": len(slide.shapes),
                    "has_images": any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes),
                    "has_text": any(hasattr(shape, "text_frame") and shape.text_frame for shape in slide.shapes)
                }
                info["slides"].append(slide_info)

            logger.info(f"PowerPointファイル情報を取得しました: {pptx_path}")
            return info

        except Exception as e:
            self.error_handler.handle_error(
                error_type=ErrorType.PPTX_GENERATION,
                error=e,
                details={"file": str(pptx_path)}
            )
            raise


# ユーティリティ関数
def estimate_pptx_size(image_paths: List[Path], settings: SlideSettings) -> float:
    """
    PowerPointファイルサイズを推定

    Args:
        image_paths: 画像ファイルパスのリスト
        settings: スライド設定

    Returns:
        推定ファイルサイズ（MB）
    """
    try:
        total_image_size = sum(path.stat().st_size for path in image_paths if path.exists())

        # 基本サイズ（空のプレゼンテーション）
        base_size = 0.5  # MB

        # 画像圧縮率を考慮（PowerPointでは通常30-50%に圧縮される）
        compressed_size = total_image_size * 0.4 / (1024 * 1024)

        # オーバーヘッド（XML構造、メタデータなど）
        overhead = len(image_paths) * 0.1  # スライド当たり0.1MB

        estimated_size = base_size + compressed_size + overhead

        return max(1.0, estimated_size)  # 最低1MB

    except Exception as e:
        logger.warning(f"ファイルサイズ推定エラー: {e}")
        return 10.0  # デフォルト値


def validate_image_for_pptx(image_path: Path) -> bool:
    """
    PowerPoint用画像ファイルの検証

    Args:
        image_path: 画像ファイルパス

    Returns:
        検証成功時True
    """
    try:
        if not image_path.exists():
            return False

        # サポートされている拡張子
        supported_extensions = {'.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff'}

        if image_path.suffix.lower() not in supported_extensions:
            return False

        # PIL で画像を開けるかテスト
        with Image.open(image_path) as img:
            img.verify()

        return True

    except Exception:
        return False