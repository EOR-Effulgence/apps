"""
シンプルな座標テスト - 実際のPowerPointファイルでの座標確認
"""

from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

def create_test_presentation():
    """座標テスト用のPowerPointファイルを作成"""

    # 新しいプレゼンテーションを作成
    prs = Presentation()

    # 空白レイアウトを使用
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    print(f"スライドサイズ: {prs.slide_width/36000:.1f} x {prs.slide_height/36000:.1f} mm")

    # テスト1: 様々な座標での四角形配置
    positions = [
        (0, 0, "座標(0,0) - 生の数値"),
        (Mm(0), Mm(0), "座標Mm(0,0) - Mm関数"),
        (Mm(0), Mm(20), "座標Mm(0,20)"),
        (Mm(10), Mm(0), "座標Mm(10,0)"),
        (Mm(10), Mm(20), "座標Mm(10,20)"),
    ]

    colors = [
        RGBColor(255, 0, 0),    # 赤
        RGBColor(0, 255, 0),    # 緑
        RGBColor(0, 0, 255),    # 青
        RGBColor(255, 255, 0),  # 黄
        RGBColor(255, 0, 255),  # マゼンタ
    ]

    for i, (left, top, description) in enumerate(positions):
        # 四角形を配置
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=left,
            top=top,
            width=Mm(60),
            height=Mm(15)
        )

        # 背景色設定
        rect.fill.solid()
        rect.fill.fore_color.rgb = colors[i]

        # テキストを追加
        textbox = slide.shapes.add_textbox(
            left=left,
            top=top,
            width=Mm(60),
            height=Mm(15)
        )

        text_frame = textbox.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = description
        p.font.name = "メイリオ"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(255, 255, 255)  # 白

        # 実際の座標値をログ出力
        print(f"配置 {i+1}: {description}")
        print(f"  指定値: left={left}, top={top}")
        print(f"  実際値: left={rect.left}, top={rect.top}")
        print(f"  EMU換算: left={rect.left/36000:.2f}mm, top={rect.top/36000:.2f}mm")
        print()

    # ファイル保存
    output_path = Path("coordinate_verification.pptx")
    prs.save(str(output_path))
    print(f"テストファイル作成: {output_path}")
    print(f"ファイルサイズ: {output_path.stat().st_size:,} bytes")

    return output_path

if __name__ == "__main__":
    create_test_presentation()