"""
PowerPointオブジェクトの座標指定テスト
python-pptxライブラリの座標系動作を確認
"""

from pptx import Presentation
from pptx.util import Inches, Mm, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def test_coordinates():
    """座標指定のテスト"""

    # 新しいプレゼンテーションを作成
    prs = Presentation()

    # 空白のスライドレイアウトを取得
    blank_slide_layout = prs.slide_layouts[6]  # 6は空白レイアウト
    slide = prs.slides.add_slide(blank_slide_layout)

    print(f"スライドサイズ: {prs.slide_width} x {prs.slide_height} EMU")
    print(f"スライドサイズ: {prs.slide_width/36000:.1f} x {prs.slide_height/36000:.1f} mm")

    # テスト1: 座標(0,0)に四角形を配置
    print("\n=== テスト1: 座標(0,0)の四角形 ===")
    rect1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=0,
        top=0,
        width=Mm(50),
        height=Mm(20)
    )
    rect1.fill.solid()
    rect1.fill.fore_color.rgb = RGBColor(255, 0, 0)  # 赤
    print(f"四角形1: left={rect1.left}, top={rect1.top}")

    # テスト2: Mm(0)を使用
    print("\n=== テスト2: Mm(0)の四角形 ===")
    rect2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Mm(0),
        top=Mm(25),  # 少し下にずらす
        width=Mm(50),
        height=Mm(20)
    )
    rect2.fill.solid()
    rect2.fill.fore_color.rgb = RGBColor(0, 255, 0)  # 緑
    print(f"四角形2: left={rect2.left}, top={rect2.top}")

    # テスト3: 少しマージンを持たせた位置
    print("\n=== テスト3: マージン付きの四角形 ===")
    rect3 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Mm(2),
        top=Mm(50),
        width=Mm(50),
        height=Mm(20)
    )
    rect3.fill.solid()
    rect3.fill.fore_color.rgb = RGBColor(0, 0, 255)  # 青
    print(f"四角形3: left={rect3.left}, top={rect3.top}")

    # テスト4: テキストボックスの座標テスト
    print("\n=== テスト4: テキストボックスの座標 ===")
    textbox1 = slide.shapes.add_textbox(
        left=0,
        top=Mm(75),
        width=Mm(80),
        height=Mm(12)
    )
    text_frame1 = textbox1.text_frame
    text_frame1.text = "座標(0,75mm)のテキスト"
    print(f"テキストボックス1: left={textbox1.left}, top={textbox1.top}")

    textbox2 = slide.shapes.add_textbox(
        left=Mm(0),
        top=Mm(90),
        width=Mm(80),
        height=Mm(12)
    )
    text_frame2 = textbox2.text_frame
    text_frame2.text = "Mm(0)のテキスト"
    print(f"テキストボックス2: left={textbox2.left}, top={textbox2.top}")

    # 各単位の変換確認
    print("\n=== 単位変換確認 ===")
    print(f"Mm(0) = {Mm(0)} EMU")
    print(f"Mm(1) = {Mm(1)} EMU")
    print(f"Mm(10) = {Mm(10)} EMU")
    print(f"Inches(0) = {Inches(0)} EMU")
    print(f"Pt(0) = {Pt(0)} EMU")

    # ファイルを保存
    output_path = "coordinate_test.pptx"
    prs.save(output_path)
    print(f"\nテストファイルを保存しました: {output_path}")

    return output_path

if __name__ == "__main__":
    test_coordinates()