"""
PowerPoint conversion service for PDF to PPTX conversion.
Provides modular, reusable PowerPoint generation functionality.
"""

from __future__ import annotations

from pathlib import Path
from typing import List, Optional, Callable
from io import BytesIO

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

from .pdf_processor import (
    ConversionConfig,
    ConversionService,
    open_pdf_document,
    process_page_to_bytes,
    mm_to_emu,
    points_to_emu,
    PDFProcessingError
)


class PowerPointConversionService(ConversionService):
    """
    Service for converting PDF files to PowerPoint presentations.
    Provides high-quality PPTX generation with configurable layout and styling.
    """

    def __init__(self, config: ConversionConfig):
        super().__init__(config)
        self.progress_callback: Optional[Callable[[int], None]] = None

    def set_progress_callback(self, callback: Callable[[int], None]) -> None:
        """Set callback function for progress updates."""
        self.progress_callback = callback

    def convert_pdf_to_powerpoint(self, pdf_path: Path, output_dir: Path) -> Path:
        """
        Convert a single PDF file to PowerPoint presentation.

        Args:
            pdf_path: Path to PDF file
            output_dir: Directory to save PPTX file

        Returns:
            Path to generated PPTX file

        Raises:
            PDFProcessingError: If conversion fails
        """
        if not pdf_path.exists() or not pdf_path.suffix.lower() == '.pdf':
            raise PDFProcessingError(f"Invalid PDF file: {pdf_path}")

        output_dir.mkdir(parents=True, exist_ok=True)

        try:
            # Create presentation with A3 landscape dimensions
            presentation = self._create_presentation()

            # Process PDF pages
            self._add_pdf_to_presentation(pdf_path, presentation)

            # Save presentation
            output_filename = f"{pdf_path.stem}.pptx"
            output_path = output_dir / output_filename
            presentation.save(str(output_path))

            return output_path

        except Exception as e:
            raise PDFProcessingError(f"Failed to convert {pdf_path} to PowerPoint: {e}")

    def convert_multiple_pdfs_to_single_presentation(
        self,
        pdf_files: List[Path],
        output_path: Path
    ) -> Path:
        """
        Convert multiple PDF files to a single PowerPoint presentation.

        Args:
            pdf_files: List of PDF file paths
            output_path: Path for output PPTX file

        Returns:
            Path to generated PPTX file

        Raises:
            PDFProcessingError: If conversion fails
        """
        if not pdf_files:
            raise PDFProcessingError("No PDF files provided for conversion")

        try:
            # Create presentation
            presentation = self._create_presentation()

            # Process each PDF file
            for pdf_file in pdf_files:
                self._add_pdf_to_presentation(pdf_file, presentation)

            # Save presentation
            output_path.parent.mkdir(parents=True, exist_ok=True)
            presentation.save(str(output_path))

            return output_path

        except Exception as e:
            raise PDFProcessingError(f"Failed to convert PDFs to PowerPoint: {e}")

    def _create_presentation(self) -> Presentation:
        """
        Create a new PowerPoint presentation with A3 landscape dimensions.

        Returns:
            Configured presentation object
        """
        presentation = Presentation()

        # Set slide dimensions to A3 landscape (420mm x 297mm)
        presentation.slide_width = mm_to_emu(self.config.slide_width_mm)
        presentation.slide_height = mm_to_emu(self.config.slide_height_mm)

        return presentation

    def _add_pdf_to_presentation(self, pdf_path: Path, presentation: Presentation) -> int:
        """
        Add all pages from a PDF file to the presentation.

        Args:
            pdf_path: PDF file to add
            presentation: Target presentation

        Returns:
            Number of slides added

        Raises:
            PDFProcessingError: If PDF processing fails
        """
        slides_added = 0
        base_name = pdf_path.stem

        try:
            with open_pdf_document(pdf_path) as doc:
                total_pages = len(doc)

                for page_num, page in enumerate(doc, start=1):
                    self._add_page_to_presentation(
                        page, presentation, base_name, page_num
                    )
                    slides_added += 1

                    # Update progress if callback is set
                    if self.progress_callback:
                        self.progress_callback(page_num)

        except Exception as e:
            raise PDFProcessingError(f"Failed to process PDF {pdf_path}: {e}")

        return slides_added

    def _add_page_to_presentation(
        self,
        page,
        presentation: Presentation,
        base_name: str,
        page_num: int
    ) -> None:
        """
        Add a single PDF page to the presentation.

        Args:
            page: PyMuPDF page object
            presentation: Target presentation
            base_name: Base filename for labeling
            page_num: Page number for labeling
        """
        # Get page dimensions and determine rotation
        rect = page.rect
        is_portrait = rect.width < rect.height
        rotation = 90 if (self.config.auto_rotate and is_portrait) else 0

        # Calculate image dimensions after rotation
        if rotation == 90:
            width_emu = points_to_emu(rect.height)
            height_emu = points_to_emu(rect.width)
        else:
            width_emu = points_to_emu(rect.width)
            height_emu = points_to_emu(rect.height)

        # Generate high-quality image
        image_bytes, _ = process_page_to_bytes(page, self.config, "png")
        image_stream = BytesIO(image_bytes)

        # Add slide with blank layout
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        # Calculate positioning to center image on slide
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        final_width, final_height = self._calculate_fitted_dimensions(
            width_emu, height_emu, slide_width, slide_height
        )

        left = int((slide_width - final_width) / 2)
        top = int((slide_height - final_height) / 2)

        # Add image to slide
        slide.shapes.add_picture(
            image_stream, left, top, width=final_width, height=final_height
        )

        # Add filename label
        self._add_filename_label(slide, f"{base_name}_page_{page_num:03d}", presentation)

    def _calculate_fitted_dimensions(
        self,
        image_width: int,
        image_height: int,
        slide_width: int,
        slide_height: int
    ) -> tuple[int, int]:
        """
        Calculate dimensions to fit image within slide while maintaining aspect ratio.

        Args:
            image_width: Original image width in EMU
            image_height: Original image height in EMU
            slide_width: Slide width in EMU
            slide_height: Slide height in EMU

        Returns:
            Tuple of (fitted_width, fitted_height) in EMU
        """
        # Leave some margin (10% on each side)
        margin_ratio = 0.9
        available_width = int(slide_width * margin_ratio)
        available_height = int(slide_height * margin_ratio)

        # If image fits within available space, use original size
        if image_width <= available_width and image_height <= available_height:
            return image_width, image_height

        # Calculate scaling ratios
        width_ratio = available_width / image_width
        height_ratio = available_height / image_height

        # Use the smaller ratio to ensure image fits
        scale_ratio = min(width_ratio, height_ratio)

        fitted_width = int(image_width * scale_ratio)
        fitted_height = int(image_height * scale_ratio)

        return fitted_width, fitted_height

    def _add_filename_label(
        self,
        slide,
        label_text: str,
        presentation: Presentation
    ) -> None:
        """
        Add filename label to the slide with configurable styling.

        Args:
            slide: Slide to add label to
            label_text: Text for the label
            presentation: Presentation for dimensions
        """
        # Import config here to avoid circular imports
        try:
            from ..config import get_app_config
            config = get_app_config()
            label_config = config.powerpoint_label
        except ImportError:
            # Fallback to default configuration
            label_config = self._get_default_label_config()

        # Calculate label dimensions
        label_width = int(presentation.slide_width * 0.25)  # 25% of slide width
        label_height = int(presentation.slide_height * 0.08)  # 8% of slide height

        # Calculate position based on configuration
        left, top = self._calculate_label_position(
            getattr(label_config, 'position', 'bottom-right'),
            presentation.slide_width,
            presentation.slide_height,
            label_width,
            label_height
        )

        # Create text box
        textbox = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left, top,
            label_width,
            label_height
        )

        # Apply styling
        self._apply_label_styling(textbox, label_config)

        # Set text properties
        text_frame = textbox.text_frame
        text_frame.margin_left = Pt(6)
        text_frame.margin_right = Pt(6)
        text_frame.margin_top = Pt(3)
        text_frame.margin_bottom = Pt(3)

        paragraph = text_frame.paragraphs[0]
        paragraph.text = label_text
        paragraph.alignment = PP_ALIGN.CENTER

        # Apply font configuration
        font_name = getattr(label_config, 'font_name', 'Arial')
        font_size = getattr(label_config, 'font_size', 12)
        font_bold = getattr(label_config, 'font_bold', False)

        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = font_bold

        # Apply text color
        text_color = getattr(label_config, 'text_color', '#000000')
        text_rgb = self._hex_to_rgb(text_color)
        paragraph.font.color.rgb = RGBColor(*text_rgb)

    def _get_default_label_config(self):
        """Get default label configuration as fallback."""
        class DefaultLabelConfig:
            position = 'bottom-right'
            text_color = '#000000'
            background_color = '#FFFFFF'
            border_color = '#FF0000'
            font_name = 'Arial'
            font_size = 12
            font_bold = False
            border_width = 1.0
            enable_shadow = False
            shadow_color = '#808080'

        return DefaultLabelConfig()

    def _calculate_label_position(
        self,
        position: str,
        slide_width: int,
        slide_height: int,
        label_width: int,
        label_height: int
    ) -> tuple[int, int]:
        """
        Calculate label position based on configuration.

        Args:
            position: Position setting
            slide_width: Slide width in EMU
            slide_height: Slide height in EMU
            label_width: Label width in EMU
            label_height: Label height in EMU

        Returns:
            Tuple of (left, top) coordinates in EMU
        """
        margin = int(min(slide_width, slide_height) * 0.02)  # 2% margin

        if position == "top-left":
            return margin, margin
        elif position == "top-right":
            return slide_width - label_width - margin, margin
        elif position == "bottom-left":
            return margin, slide_height - label_height - margin
        elif position == "bottom-right":
            return slide_width - label_width - margin, slide_height - label_height - margin
        else:
            # Default to bottom-right
            return slide_width - label_width - margin, slide_height - label_height - margin

    def _apply_label_styling(self, textbox, label_config) -> None:
        """
        Apply styling configuration to label textbox.

        Args:
            textbox: PowerPoint textbox shape
            label_config: Label configuration object
        """
        # Apply background color
        textbox.fill.solid()
        bg_color = getattr(label_config, 'background_color', '#FFFFFF')
        bg_rgb = self._hex_to_rgb(bg_color)
        textbox.fill.fore_color.rgb = RGBColor(*bg_rgb)

        # Apply border styling
        border_color = getattr(label_config, 'border_color', '#FF0000')
        border_width = getattr(label_config, 'border_width', 1.0)
        border_rgb = self._hex_to_rgb(border_color)
        textbox.line.color.rgb = RGBColor(*border_rgb)
        textbox.line.width = int(border_width * 12700)  # Convert to EMU

        # Apply shadow if enabled
        enable_shadow = getattr(label_config, 'enable_shadow', False)
        if enable_shadow:
            textbox.shadow.inherit = False
            textbox.shadow.visible = True
            shadow_color = getattr(label_config, 'shadow_color', '#808080')
            shadow_rgb = self._hex_to_rgb(shadow_color)
            textbox.shadow.color.rgb = RGBColor(*shadow_rgb)
        else:
            textbox.shadow.inherit = False

    def _hex_to_rgb(self, hex_color: str) -> tuple[int, int, int]:
        """
        Convert hex color to RGB tuple.

        Args:
            hex_color: Hex color string (e.g., '#FF0000')

        Returns:
            RGB tuple (r, g, b)
        """
        hex_color = hex_color.lstrip('#')
        if len(hex_color) != 6:
            return (0, 0, 0)  # Default to black if invalid

        try:
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except ValueError:
            return (0, 0, 0)  # Default to black if invalid

    def convert(self, input_path: Path, output_path: Path) -> bool:
        """
        Convert single PDF to PowerPoint (implementation of abstract method).

        Args:
            input_path: PDF file path
            output_path: Output directory path

        Returns:
            True if conversion successful

        Raises:
            PDFProcessingError: If conversion fails
        """
        output_file = self.convert_pdf_to_powerpoint(input_path, output_path)
        return output_file.exists()